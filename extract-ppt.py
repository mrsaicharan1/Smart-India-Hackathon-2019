from pptx import Presentation
import os

cur_directory = os.path.dirname(os.path.abspath( __file__ ))
files = [x for x in os.listdir(cur_directory) if x.endswith(".pptx")]


for eachfile in files:
    prs = Presentation(eachfile)
    print(eachfile)
    print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)